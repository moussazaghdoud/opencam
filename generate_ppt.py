from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG = RGBColor(0x0A, 0x0A, 0x0A)
CARD_BG = RGBColor(0x14, 0x14, 0x14)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xED, 0xED, 0xED)
GRAY = RGBColor(0x73, 0x73, 0x73)
LGRAY = RGBColor(0xA0, 0xA0, 0xA0)
DCARD = RGBColor(0x1E, 0x1E, 0x1E)
BORDER = RGBColor(0x26, 0x26, 0x26)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def box(slide, l, t, w, h, fill=CARD_BG, border=BORDER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb


def line(slide, l, t, w, color=BLUE):
    box(slide, l, t, w, Pt(3), color, None)


# ---- SLIDE 1: Title ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(1), Inches(1.5), Inches(11), Inches(1.5), "OpenCam", 60, BLUE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3), Inches(11), Inches(1), "AI-Powered Video Surveillance Platform", 28, WHITE, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4), Inches(11), Inches(0.8), "Privacy-First  |  On-Premise AI  |  Enterprise Grade", 18, GRAY, False, PP_ALIGN.CENTER)
line(s, Inches(4), Inches(5.5), Inches(5))
txt(s, Inches(1), Inches(6.2), Inches(11), Inches(0.5), "Confidential - SecureLLM Technologies", 12, GRAY, False, PP_ALIGN.CENTER)

# ---- SLIDE 2: The Problem ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "The Problem", 36, WHITE, True)
line(s, Inches(0.8), Inches(1.2), Inches(2))

problems = [
    ("Cameras record, nobody watches", "Footage reviewed only AFTER incidents. Billions of hours of video wasted."),
    ("Cloud surveillance = privacy risk", "Video sent to US/China servers. GDPR non-compliant. Data sovereignty lost."),
    ("Existing systems are expensive", "Enterprise VMS costs 50-200K EUR. Out of reach for most businesses."),
    ("No actionable intelligence", "Raw video, no analytics, no real-time alerts, no operational insights."),
]
for i, (title, desc) in enumerate(problems):
    y = Inches(1.8 + i * 1.3)
    box(s, Inches(0.8), y, Inches(11.5), Inches(1.1))
    txt(s, Inches(1.2), y + Inches(0.15), Inches(10), Inches(0.4), title, 20, RED, True)
    txt(s, Inches(1.2), y + Inches(0.55), Inches(10), Inches(0.4), desc, 15, LGRAY)

# ---- SLIDE 3: Our Solution ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "Our Solution: OpenCam", 36, WHITE, True)
line(s, Inches(0.8), Inches(1.2), Inches(2))
txt(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8), "Turn your existing cameras into an intelligent security system.\nAI processes locally. Video never leaves your building.", 18, LGRAY)

features = [
    ("Real-Time Detection", "YOLO v8 - persons, vehicles,\nobjects detected instantly", BLUE),
    ("Face Recognition", "ArcFace - 99%+ accuracy,\nidentify or flag unknowns", GREEN),
    ("PPE Compliance", "Safety jacket detection,\nauto-alerts on violations", YELLOW),
    ("People Counting", "Directional IN/OUT counting,\nthroughput analytics", BLUE),
]
for i, (title, desc, color) in enumerate(features):
    x = Inches(0.8 + i * 3.05)
    box(s, x, Inches(3), Inches(2.8), Inches(2.5))
    box(s, x, Inches(3), Inches(2.8), Pt(3), color, None)
    txt(s, x + Inches(0.2), Inches(3.3), Inches(2.4), Inches(0.5), title, 18, color, True)
    txt(s, x + Inches(0.2), Inches(3.9), Inches(2.4), Inches(1.2), desc, 14, LGRAY)

# ---- SLIDE 4: Architecture ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "Hybrid Architecture - Privacy by Design", 36, WHITE, True)
line(s, Inches(0.8), Inches(1.2), Inches(2))

box(s, Inches(0.8), Inches(2), Inches(5), Inches(4.5), CARD_BG, GREEN)
txt(s, Inches(1), Inches(2.1), Inches(4.5), Inches(0.5), "ON-PREMISE (Customer Site)", 16, GREEN, True)
items = ["IP Cameras (RTSP/ONVIF)", "OpenCam AI Engine", "YOLO v8 Object Detection", "ArcFace Face Recognition", "PPE / Safety Monitoring", "People Counting Engine", "Video stays HERE"]
for i, item in enumerate(items):
    c = GREEN if i == 6 else WHITE
    txt(s, Inches(1.3), Inches(2.7 + i * 0.45), Inches(4), Inches(0.4), "> " + item, 13, c)

box(s, Inches(5.8), Inches(3.6), Inches(1.5), Inches(1.2), DCARD, BLUE)
txt(s, Inches(5.9), Inches(3.8), Inches(1.3), Inches(0.8), "NATS Tunnel\n(encrypted)", 11, BLUE, False, PP_ALIGN.CENTER)

box(s, Inches(7.5), Inches(2), Inches(5), Inches(4.5), CARD_BG, BLUE)
txt(s, Inches(7.7), Inches(2.1), Inches(4.5), Inches(0.5), "CLOUD (SaaS Dashboard)", 16, BLUE, True)
items2 = ["Multi-site Management", "Real-time Dashboard", "Alert Hub (email, Slack, webhook)", "User Management / SSO", "Analytics & Reports", "LLM Insights (Claude API)", "Only metadata - no video"]
for i, item in enumerate(items2):
    c = BLUE if i == 6 else WHITE
    txt(s, Inches(7.9), Inches(2.7 + i * 0.45), Inches(4), Inches(0.4), "> " + item, 13, c)

# ---- SLIDE 5: Technology Stack ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "Enterprise-Grade Technology Stack", 36, WHITE, True)
line(s, Inches(0.8), Inches(1.2), Inches(2))

protocols = [
    ("RTSP / ONVIF", "Camera streaming", "Universal industry standard\n20,000+ compatible cameras"),
    ("H.264 / H.265", "Video codec", "Mandatory in all\nprofessional surveillance"),
    ("YOLOv8", "Object detection", "State-of-the-art accuracy\nUsed by Avigilon, BriefCam"),
    ("ArcFace (ONNX)", "Face recognition", "99.8% accuracy benchmark\nUsed by NEC, law enforcement"),
    ("NATS", "Cloud messaging", "Enterprise-grade\nUsed by VMware, Mastercard"),
    ("FastAPI + PostgreSQL", "Backend", "Production-grade stack\nUsed by Netflix, Uber"),
]
for i, (name, cat, desc) in enumerate(protocols):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.8 + row * 2.5)
    box(s, x, y, Inches(3.8), Inches(2.1))
    txt(s, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.4), name, 20, BLUE, True)
    txt(s, x + Inches(0.2), y + Inches(0.55), Inches(3.4), Inches(0.3), cat, 12, GRAY)
    txt(s, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(0.8), desc, 13, LGRAY)
    txt(s, x + Inches(3.0), y + Inches(0.15), Inches(0.5), Inches(0.4), "PRO", 11, GREEN, True)

# ---- SLIDE 6: Compliance ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "Compliance & Security", 36, WHITE, True)
line(s, Inches(0.8), Inches(1.2), Inches(2))

stds = [
    ("GDPR Compliant", "Video processed on-premise\nNo data leaves the building\nPrivacy by design", GREEN),
    ("ONVIF Certified", "Profile S & T compatible\nCamera interoperability\n20,000+ devices supported", BLUE),
    ("NDAA Compliant", "Compatible cameras available\nNo banned manufacturers\nGovernment-ready", BLUE),
    ("ISO 27001 Ready", "Architecture supports\ndata security requirements\nEncrypted communications", GREEN),
]
for i, (title, desc, color) in enumerate(stds):
    x = Inches(0.8 + i * 3.05)
    box(s, x, Inches(2), Inches(2.8), Inches(3))
    box(s, x, Inches(2), Inches(2.8), Pt(3), color, None)
    txt(s, x + Inches(0.2), Inches(2.4), Inches(2.4), Inches(0.5), title, 20, color, True)
    txt(s, x + Inches(0.2), Inches(3.1), Inches(2.4), Inches(1.5), desc, 14, LGRAY)

# ---- SLIDE 7: Competitive Advantage ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "Competitive Advantage", 36, WHITE, True)
line(s, Inches(0.8), Inches(1.2), Inches(2))

box(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6), BLUE, None)
headers = ["Feature", "Basic CCTV", "Cloud (Verkada)", "OpenCam"]
widths = [3, 2.5, 3, 3]
xp = 0.8
for h, w in zip(headers, widths):
    txt(s, Inches(xp), Inches(1.85), Inches(w), Inches(0.5), h, 14, WHITE, True, PP_ALIGN.CENTER)
    xp += w

rows = [
    ("AI Detection", "No", "Yes (cloud)", "Yes (local)"),
    ("Face Recognition", "No", "Limited", "ArcFace 99%+"),
    ("Video Privacy", "Local only", "Sent to US cloud", "Stays on-premise"),
    ("GDPR Compliant", "Yes", "Questionable", "Yes - by design"),
    ("Setup Cost", "5-20K EUR", "50-200K EUR", "Under 500 EUR"),
    ("Monthly Cost", "0 (no analytics)", "500-5000 EUR", "99-999 EUR"),
    ("Real-time Alerts", "No", "Yes", "Yes + customizable"),
    ("Works Offline", "Yes", "No", "Yes"),
]
for i, (feat, cctv, cloud, oc) in enumerate(rows):
    y = Inches(2.5 + i * 0.55)
    bg = CARD_BG if i % 2 == 0 else DCARD
    box(s, Inches(0.8), y, Inches(11.5), Inches(0.5), bg, None)
    xp = 0.8
    vals = [feat, cctv, cloud, oc]
    colors = [WHITE, RED if cctv in ("No", "Questionable") else LGRAY,
              YELLOW if "cloud" in cloud.lower() or cloud == "Questionable" else LGRAY, GREEN]
    for v, w, c in zip(vals, widths, colors):
        txt(s, Inches(xp), y + Inches(0.05), Inches(w), Inches(0.4), v, 12, c, False, PP_ALIGN.CENTER)
        xp += w

# ---- SLIDE 8: Target Markets ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "Target Markets", 36, WHITE, True)
line(s, Inches(0.8), Inches(1.2), Inches(2))

markets = [
    ("Logistics & Warehouses", "People counting, dock timing\nPPE compliance, throughput KPIs", "50-200K EUR/yr saved"),
    ("Banking & Finance", "Access control, face ID\nVault monitoring, compliance", "Mandatory by regulation"),
    ("Retail Chains", "Customer counting, loss prevention\nStaff optimization, analytics", "10-30% shrinkage reduction"),
    ("Healthcare", "Patient safety, restricted zones\nStaff tracking, compliance", "Critical safety requirement"),
    ("Industrial Sites", "PPE enforcement, zone control\nForklift safety, incident detection", "100K+ EUR fines avoided"),
    ("Government & Public", "Building security, face ID\nCrowd management, NDAA compliant", "Security mandate"),
]
for i, (title, desc, impact) in enumerate(markets):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.8 + row * 2.7)
    box(s, x, y, Inches(3.8), Inches(2.3))
    txt(s, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.4), title, 16, BLUE, True)
    txt(s, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(1), desc, 12, LGRAY)
    box(s, x + Inches(0.2), y + Inches(1.7), Inches(3.4), Inches(0.4), DCARD, None)
    txt(s, x + Inches(0.3), y + Inches(1.75), Inches(3.2), Inches(0.3), impact, 11, GREEN, True)

# ---- SLIDE 9: Pricing ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "Pricing", 36, WHITE, True)
line(s, Inches(0.8), Inches(1.2), Inches(2))

tiers = [
    ("Starter", "99 EUR/mo", "Up to 8 cameras\n1 site\nEmail alerts\nBasic analytics\nFace recognition", LGRAY),
    ("Business", "299 EUR/mo", "Up to 32 cameras\n5 sites\nAll alert channels\nFull analytics\nPPE detection\nAPI access", BLUE),
    ("Enterprise", "999+ EUR/mo", "Unlimited cameras\nUnlimited sites\nSSO / SAML\nCustom integrations\nSLA & support\nLLM insights", GREEN),
]
for i, (name, price, feats, color) in enumerate(tiers):
    x = Inches(1.3 + i * 3.8)
    box(s, x, Inches(1.8), Inches(3.3), Inches(5), CARD_BG, color)
    box(s, x, Inches(1.8), Inches(3.3), Pt(3), color, None)
    txt(s, x + Inches(0.2), Inches(2.1), Inches(2.9), Inches(0.5), name, 22, color, True, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(2.7), Inches(2.9), Inches(0.6), price, 28, WHITE, True, PP_ALIGN.CENTER)
    txt(s, x + Inches(0.3), Inches(3.6), Inches(2.7), Inches(3), feats, 14, LGRAY)

txt(s, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
    "Hardware appliance: 500 EUR one-time per site (optional)", 12, GRAY, False, PP_ALIGN.CENTER)

# ---- SLIDE 10: Contact ----
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(1), Inches(2), Inches(11), Inches(1.5), "Ready to secure\nyour operations?", 44, WHITE, True, PP_ALIGN.CENTER)
line(s, Inches(4), Inches(4), Inches(5))
txt(s, Inches(1), Inches(4.5), Inches(11), Inches(0.6), "Live demo available - see it in action today", 20, BLUE, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "Moussa Zaghdoud - moussa.zaghdoud@gmail.com", 16, LGRAY, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.2), Inches(11), Inches(0.5), "SecureLLM Technologies", 14, GRAY, False, PP_ALIGN.CENTER)

prs.save("C:/Users/zaghdoud/opencam/OpenCam_Presentation.pptx")
print("Done! Saved to OpenCam_Presentation.pptx")
